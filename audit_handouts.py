# audit_handouts.py
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, TypedDict, Union
import io
import re
import json
import hashlib
import time

import requests
from docxtpl import DocxTemplate
from docx import Document  # python-docx


class HandoutChunk(TypedDict):
    file: str
    unit: str
    heading: str
    text: str
    captions: List[str]
    category: str  # "lecture" | "lab"

_CAPTION_RE = re.compile(
    r"(?im)^\s*(figure|fig\.|table|eq\.|equation)\s*\d+"
    r"(\s*[:.\-–—]\s*.*)?$"
)

def _sha256(b: bytes) -> str:
    return hashlib.sha256(b).hexdigest()

def _first_nonempty_line(text: str) -> str:
    for line in (text or "").splitlines():
        s = line.strip()
        if s:
            return s[:200]
    return ""

def _extract_captions(text: str) -> List[str]:
    if not text:
        return []
    caps = []
    for m in _CAPTION_RE.finditer(text):
        line = (m.group(0) or "").strip()
        if line and line not in caps:
            caps.append(line[:250])
    return caps

def _clean_text(s: str) -> str:
    if not s:
        return ""
    # normalize whitespace without destroying structure too much
    s = s.replace("\u00ad", "")  # soft hyphen
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()

def parse_pdf(path: Union[str, Path]) -> List[HandoutChunk]:
    path = Path(path)
    # Prefer pymupdf (fitz), else pypdf.
    try:
        import fitz  # pymupdf
        doc = fitz.open(str(path))
        out: List[HandoutChunk] = []
        for i in range(doc.page_count):
            page = doc.load_page(i)
            text = page.get_text("text") or ""
            text = _clean_text(text)
            out.append({
                "file": path.name,
                "unit": f"page {i+1}",
                "heading": _first_nonempty_line(text),
                "text": text,
                "captions": _extract_captions(text),
            })
        return out
    except Exception:
        pass

    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        out: List[HandoutChunk] = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            text = _clean_text(text)
            out.append({
                "file": path.name,
                "unit": f"page {i+1}",
                "heading": _first_nonempty_line(text),
                "text": text,
                "captions": _extract_captions(text),
            })
        return out
    except Exception as e:
        raise RuntimeError(f"PDF parsing failed (install pymupdf or pypdf). Root error: {e}") from e

def parse_pptx(path: Union[str, Path]) -> List[HandoutChunk]:
    path = Path(path)
    try:
        from pptx import Presentation
    except Exception as e:
        raise RuntimeError("python-pptx is required for PPTX parsing. Please `pip install python-pptx`.") from e

    prs = Presentation(str(path))
    out: List[HandoutChunk] = []

    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        try:
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
        except Exception:
            title = ""

        texts: List[str] = []
        for shape in slide.shapes:
            try:
                if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                    t = shape.text_frame.text
                    if t and t.strip():
                        texts.append(t.strip())
            except Exception:
                continue

        notes = ""
        try:
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            notes = ""

        full = "\n".join(texts)
        if notes:
            full = (full + "\n\n[NOTES]\n" + notes).strip()

        full = _clean_text(full)
        heading = title.strip() if title.strip() else _first_nonempty_line(full)

        out.append({
            "file": path.name,
            "unit": f"slide {idx}",
            "heading": heading,
            "text": full,
            "captions": _extract_captions(full),
        })

    return out

def build_corpus(paths: List[Union[str, Path]]) -> List[HandoutChunk]:
    """Parse + normalize into a single ordered corpus."""
    corpus: List[HandoutChunk] = []
    for p in paths:
        p = Path(p)
        ext = p.suffix.lower()
        if ext == ".pdf":
            corpus.extend(parse_pdf(p))
        elif ext == ".pptx":
            corpus.extend(parse_pptx(p))
        elif ext == ".ppt":
            # .ppt is not supported by python-pptx
            corpus.append({
                "file": p.name,
                "unit": "file",
                "heading": "unsupported .ppt",
                "text": "Unsupported format: .ppt (please export/save as .pptx).",
                "captions": [],
            })
        else:
            corpus.append({
                "file": p.name,
                "unit": "file",
                "heading": "unsupported",
                "text": f"Unsupported format: {ext}",
                "captions": [],
            })
    return corpus

def _compact_chunk(c: HandoutChunk, max_chars: int = 2600) -> HandoutChunk:
    t = c.get("text", "") or ""
    if len(t) <= max_chars:
        return c
    # keep start + end, preserve some structure
    head = t[: int(max_chars * 0.75)]
    tail = t[- int(max_chars * 0.2):]
    t2 = (head + "\n...\n" + tail).strip()
    out: HandoutChunk = {
        "file": c["file"],
        "unit": c["unit"],
        "heading": c.get("heading", "") or "",
        "text": t2,
        "captions": (c.get("captions") or [])[:15],
    }
    return out

def _flatten_weekly(draft_bundle: Dict[str, Any]) -> str:
    w = (draft_bundle or {}).get("weekly", {}) or {}
    theory = w.get("theory_rows", []) or []
    practical = w.get("practical_rows", []) or []
    def fmt(kind: str, rows: List[Dict[str, Any]]) -> List[str]:
        out = [f"[{kind}]"]
        for i, r in enumerate(rows, start=1):
            if not any((r.get(k) for k in ["week","topic","hours","methods","assessment","clos","gas"])):
                continue
            out.append(
                f"- Row {i}: week={r.get('week','')}, hours={r.get('hours','')}, "
                f"topic={r.get('topic','')}, CLOs={r.get('clos','')}, GAs={r.get('gas','')}, "
                f"methods={r.get('methods','')}, assessment={r.get('assessment','')}"
            )
        return out
    lines = fmt("Theory", theory) + [""] + fmt("Practical", practical)
    return "\n".join(lines).strip()

AUDIT_SYSTEM_PROMPT = """You are an academic course handout auditor for UTAS.
You must output ONLY valid JSON matching the exact contract provided by the user.
Ground truth:
- Use the CDP snapshot (goals/CLOs/weekly topics/sources) as ground truth.
Evidence rules:
- Use ONLY the uploaded handouts as evidence; be conservative.
- Handouts may include:
  (A) current lecture handouts
  (B) optional lab/practical experiment handouts
  (C) optional previous semester versions of either/both
- If you cannot confidently support a judgement, set evidence="insufficient" and explain what is missing in remarks.
- When lab handouts are provided, treat them as evidence for the PRACTICAL portion of the CDP weekly distribution (practical rows, practical CLO/GA coverage, practical methods/assessment alignment).
Previous comparison:
- If previous handouts are provided, comment ONLY on observable differences; otherwise mark updated_vs_previous as not_applicable.
Output constraints:
- Never include markdown fences or extra commentary outside JSON.
- Keep evidence specific: mention file + unit (page/slide) and a short quote/summary pointer.
"""


def prepare_llm_payload(
    cdp_bundle: Dict[str, Any],
    corpus: List[HandoutChunk],
    corpus_prev: Optional[List[HandoutChunk]] = None,
    corpus_lab_current: Optional[List[HandoutChunk]] = None,
    corpus_lab_previous: Optional[List[HandoutChunk]] = None,
) -> Dict[str, Any]:
    course = (cdp_bundle or {}).get("course", {}) or {}
    doc    = (cdp_bundle or {}).get("doc", {}) or {}
    semester = str(course.get("semester") or doc.get("semester") or "").strip()
    academic_year = str(course.get("academic_year") or doc.get("academic_year") or "").strip()

    def _tag(corp: Optional[List[HandoutChunk]], cat: str) -> Optional[List[HandoutChunk]]:
        if not corp:
            return None
        out: List[HandoutChunk] = []
        for c in corp:
            cc = dict(c)
            cc["category"] = cat
            out.append(cc)  # type: ignore
        return out

    def _cap(xs: List[HandoutChunk], max_n: int) -> List[HandoutChunk]:
        if len(xs) <= max_n:
            return xs
        head_n = max_n // 2
        tail_n = max_n - head_n
        return xs[:head_n] + xs[-tail_n:]

    # --- infer material flags for doc rendering (optional but useful) ---
    def _has_ext(chunks: Optional[List[HandoutChunk]], exts: Tuple[str, ...]) -> bool:
        if not chunks:
            return False
        for c in chunks:
            fn = (c.get("file") or "").lower()
            if any(fn.endswith(e) for e in exts):
                return True
        return False

    cdp_snapshot = {
        "course_code": course.get("course_code", ""),
        "course_title": course.get("course_title", ""),
        "semester": semester,
        "academic_year": academic_year,
        "goals": (cdp_bundle or {}).get("goals_text", "") or "",
        "clos": (cdp_bundle or {}).get("clos_table", "") or "",
        "weekly_topics": _flatten_weekly(cdp_bundle),
        "sources": (cdp_bundle or {}).get("sources_text", "") or "",
        "materials": {
            "has_ppt": _has_ext(corpus, (".pptx", ".ppt")),
            "has_pdf": _has_ext(corpus, (".pdf",)),
            "has_lab_manual": bool(corpus_lab_current),
        },
    }

    # Tag corpora (category helps the model interpret them)
    corpus_lecture_cur = _tag(corpus, "lecture") or []
    corpus_lecture_prev = _tag(corpus_prev, "lecture") if corpus_prev else None
    corpus_lab_cur = _tag(corpus_lab_current, "lab") if corpus_lab_current else None
    corpus_lab_prev = _tag(corpus_lab_previous, "lab") if corpus_lab_previous else None

    current_compact = [_compact_chunk(c) for c in corpus_lecture_cur if (c.get("text") or "").strip()]
    current_compact = _cap(current_compact, 70)

    prev_compact = None
    if corpus_lecture_prev:
        prev_compact = [_compact_chunk(c) for c in corpus_lecture_prev if (c.get("text") or "").strip()]
        prev_compact = _cap(prev_compact, 50)

    lab_current_compact = None
    if corpus_lab_cur:
        lab_current_compact = [_compact_chunk(c) for c in corpus_lab_cur if (c.get("text") or "").strip()]
        lab_current_compact = _cap(lab_current_compact, 70)

    lab_prev_compact = None
    if corpus_lab_prev:
        lab_prev_compact = [_compact_chunk(c) for c in corpus_lab_prev if (c.get("text") or "").strip()]
        lab_prev_compact = _cap(lab_prev_compact, 50)

    user_payload = {
        "cdp_snapshot": cdp_snapshot,

        # lecture handouts (existing)
        "handouts_current": current_compact,
        "handouts_previous": prev_compact,

        # NEW: lab handouts (optional)
        "handouts_lab_current": lab_current_compact,
        "handouts_lab_previous": lab_prev_compact,

        "contract": {
            "pc_review": [
                {"criterion": "coverage_vs_outcomes", "rating": "below|meet|above", "evidence": "...", "remarks": "..."},
                {"criterion": "updated_vs_previous",  "rating": "below|meet|above|not_applicable", "evidence": "...", "remarks": "..."},
                {"criterion": "innovative_methods",   "rating": "below|meet|above", "evidence": "...", "remarks": "..."},
                {"criterion": "logical_sequence",     "rating": "below|meet|above", "evidence": "...", "remarks": "..."},
            ],
            "cc_review": [
                {"criterion": "template_utilized",      "yes_no": "yes|no", "evidence": "...", "remarks": "..."},
                {"criterion": "outcomes_per_chapter",   "yes_no": "yes|no", "evidence": "...", "remarks": "..."},
                {"criterion": "numbering_structure",    "yes_no": "yes|no", "evidence": "...", "remarks": "..."},
                {"criterion": "linguistic_clarity",     "yes_no": "yes|no", "evidence": "...", "remarks": "..."},
                {"criterion": "organization",           "yes_no": "yes|no", "evidence": "...", "remarks": "..."},
                {"criterion": "captions_present",       "yes_no": "yes|no", "evidence": "...", "remarks": "..."},
                {"criterion": "proper_citation",        "yes_no": "yes|no", "evidence": "...", "remarks": "..."},
                {"criterion": "relevant_references",    "yes_no": "yes|no", "evidence": "...", "remarks": "..."},
            ],
            "overall_summary": "...",
            "action_items": ["..."],
            "trace": [
                {"chunk": {"file": "...", "unit": "slide 7", "heading": "..."}, "highlights": ["..."], "linked_criteria": ["coverage_vs_outcomes"]}
            ],
        }
    }

    user_prompt = (
        "Evaluate the course handouts against the CDP snapshot.\n"
        "Return ONLY JSON that matches the contract inside this payload.\n"
        "Use BOTH lecture handouts and (if provided) lab/practical experiment handouts as evidence.\n\n"
        + json.dumps(user_payload, ensure_ascii=False)
    )

    return {
        "system": AUDIT_SYSTEM_PROMPT,
        "user": user_prompt,
        "cdp_snapshot": cdp_snapshot,
    }

def _extract_json(text: str) -> Dict[str, Any]:
    # Claude is usually clean, but guard anyway.
    s = (text or "").strip()
    if not s:
        raise ValueError("Empty model response.")
    if s[0] != "{":
        i = s.find("{")
        j = s.rfind("}")
        if i >= 0 and j > i:
            s = s[i:j+1]
    return json.loads(s)

def run_handout_audit(
    payload: Dict[str, Any],
    api_key: str,
    model: str = "anthropic/claude-3.5-sonnet",
    app_url: Optional[str] = None,
    app_title: str = "UTAS CDP Builder",
    timeout_s: int = 120,
) -> Dict[str, Any]:
    if not api_key:
        raise RuntimeError("Missing OpenRouter API key.")

    body = {
        "model": model,
        "messages": [
            {"role": "system", "content": payload["system"]},
            {"role": "user", "content": payload["user"]},
        ],
        "temperature": 0.1,
        "max_tokens": 2400,
    }

    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }
    if app_url:
        headers["HTTP-Referer"] = app_url
    if app_title:
        headers["X-Title"] = app_title

    r = requests.post(
        "https://openrouter.ai/api/v1/chat/completions",
        headers=headers,
        json=body,
        timeout=timeout_s,
    )
    if r.status_code >= 400:
        raise RuntimeError(f"OpenRouter error {r.status_code}: {r.text[:4000]}")

    data = r.json()
    txt = (
        data.get("choices", [{}])[0]
            .get("message", {})
            .get("content", "")
    )
    out = _extract_json(txt)

    # Minimal normalization/guardrails
    out.setdefault("pc_review", [])
    out.setdefault("cc_review", [])
    out.setdefault("overall_summary", "")
    out.setdefault("action_items", [])
    out.setdefault("trace", [])
    return out

def ensure_audit_summary_template(template_path: Union[str, Path]) -> Path:
    template_path = Path(template_path)
    template_path.parent.mkdir(parents=True, exist_ok=True)
    if template_path.exists():
        return template_path

    doc = Document()
    doc.add_heading("Course Handout Audit Summary", level=1)
    doc.add_paragraph("Course: {{ course_code }} — {{ course_title }}")
    doc.add_paragraph("Academic Year: {{ academic_year }}   |   Semester: {{ semester }}")
    doc.add_paragraph("Generated at: {{ generated_at }}")

    doc.add_heading("Program Coordinator Review", level=2)
    doc.add_paragraph("{% for r in pc_review %}")
    doc.add_paragraph("• {{ r.criterion }} — {{ r.rating }}")
    doc.add_paragraph("  Evidence: {{ r.evidence }}")
    doc.add_paragraph("  Remarks: {{ r.remarks }}")
    doc.add_paragraph("{% endfor %}")

    doc.add_heading("Curriculum Committee Review", level=2)
    doc.add_paragraph("{% for r in cc_review %}")
    doc.add_paragraph("• {{ r.criterion }} — {{ r.yes_no }}")
    doc.add_paragraph("  Evidence: {{ r.evidence }}")
    doc.add_paragraph("  Remarks: {{ r.remarks }}")
    doc.add_paragraph("{% endfor %}")

    doc.add_heading("Overall Summary", level=2)
    doc.add_paragraph("{{ overall_summary }}")

    doc.add_heading("Action Items", level=2)
    doc.add_paragraph("{% for a in action_items %}• {{ a }}{% endfor %}")

    doc.save(str(template_path))
    return template_path

def render_audit_summary_docx(
    template_path: Union[str, Path],
    audit_json: Dict[str, Any],
    cdp_snapshot: Dict[str, Any],
) -> bytes:
    template_path = ensure_audit_summary_template(template_path)
    tpl = DocxTemplate(str(template_path))

    ctx = {
        "course_code": cdp_snapshot.get("course_code", ""),
        "lecturer_name": cdp_snapshot.get("lecturer_name"),
        "course_title": cdp_snapshot.get("course_title", ""),
        "academic_year": cdp_snapshot.get("academic_year", ""),
        "semester": cdp_snapshot.get("semester", ""),
        "generated_at": time.strftime("%Y-%m-%d %H:%M"),
        "pc_review": audit_json.get("pc_review", []) or [],
        "cc_review": audit_json.get("cc_review", []) or [],
        "overall_summary": audit_json.get("overall_summary", "") or "",
        "action_items": audit_json.get("action_items", []) or [],
    }
    tpl.render(ctx)

    buf = io.BytesIO()
    tpl.save(buf)
    return buf.getvalue()

def render_course_audit_form_docx(
    template_path: Union[str, Path],
    audit_json: Dict[str, Any],
    cdp_snapshot: Dict[str, Any],
    specialization: str = "",
    unit: str = "",  # "CAE" / "EECE" / "MCE"
    pc_name_sign: str = "",
    cc_member_name_sign: str = "",
    staff_ack_name_sign: str = "",
) -> bytes:
    template_path = Path(template_path)
    if not template_path.exists():
        raise FileNotFoundError(f"Course Audit Form template not found: {template_path}")

    def tick(x: bool) -> str:
        return "☑" if x else ""

    def norm(s: Any) -> str:
        return str(s or "").strip().lower()

    def find_pc(criteria_contains: str) -> Dict[str, Any]:
        for r in (audit_json.get("pc_review") or []):
            if criteria_contains in norm(r.get("criterion")):
                return r
        return {}

    def find_cc(criteria_contains: str) -> Dict[str, Any]:
        for r in (audit_json.get("cc_review") or []):
            if criteria_contains in norm(r.get("criterion")):
                return r
        return {}

    def rating_ticks(rating: str):
        rt = norm(rating)
        return (
            tick("below" in rt),
            tick("meet" in rt),
            tick("above" in rt),
        )

    def yn_ticks(yes_no: str):
        y = norm(yes_no)
        is_yes = y.startswith("y") or y == "yes" or y == "true"
        is_no  = y.startswith("n") or y == "no"  or y == "false"
        return tick(is_yes), tick(is_no)

    # ---- Basic header fields (adjust keys if your snapshot uses different names)
    course_code  = cdp_snapshot.get("course_code", "")
    course_title = cdp_snapshot.get("course_title", "")
    academic_year = cdp_snapshot.get("academic_year", "")
    semester = cdp_snapshot.get("semester", "")

    # ---- Material type inference (simple heuristic)
    # If you already store filenames in session_state, pass them in and replace this logic.
    materials = (cdp_snapshot or {}).get("materials", {}) or {}
    material_ppt = "☑" if materials.get("has_ppt") else ""
    material_handout = "☑" if (materials.get("has_pdf") or materials.get("has_ppt")) else ""
    material_reference = ""
    material_lab_manual = "☑" if materials.get("has_lab_manual") else ""


    # ---- PC rows
    pc_cov = find_pc("coverage")
    pc_upd = find_pc("updated")
    pc_inn = find_pc("innov")
    pc_seq = find_pc("logical") or find_pc("sequence")

    pc_cov_b, pc_cov_m, pc_cov_a = rating_ticks(pc_cov.get("rating"))
    pc_upd_b, pc_upd_m, pc_upd_a = rating_ticks(pc_upd.get("rating"))
    pc_inn_b, pc_inn_m, pc_inn_a = rating_ticks(pc_inn.get("rating"))
    pc_seq_b, pc_seq_m, pc_seq_a = rating_ticks(pc_seq.get("rating"))

    def remarks_blob(r: Dict[str, Any]) -> str:
        ev = (r.get("evidence") or "").strip()
        rm = (r.get("remarks") or "").strip()
        if ev and rm:
            return f"{rm}\nEvidence: {ev}"
        return rm or (f"Evidence: {ev}" if ev else "")

    # ---- CC rows
    cc_tpl = find_cc("template")
    cc_out = find_cc("outcome")
    cc_num = find_cc("number")
    cc_lng = find_cc("lingu") or find_cc("clarity")
    cc_org = find_cc("organ")
    cc_cap = find_cc("caption")
    cc_cit = find_cc("citat")
    cc_ref = find_cc("reference")

    cc_tpl_y, cc_tpl_n = yn_ticks(cc_tpl.get("yes_no"))
    cc_out_y, cc_out_n = yn_ticks(cc_out.get("yes_no"))
    cc_num_y, cc_num_n = yn_ticks(cc_num.get("yes_no"))
    cc_lng_y, cc_lng_n = yn_ticks(cc_lng.get("yes_no"))
    cc_org_y, cc_org_n = yn_ticks(cc_org.get("yes_no"))
    cc_cap_y, cc_cap_n = yn_ticks(cc_cap.get("yes_no"))
    cc_cit_y, cc_cit_n = yn_ticks(cc_cit.get("yes_no"))
    cc_ref_y, cc_ref_n = yn_ticks(cc_ref.get("yes_no"))

    overall = (audit_json.get("overall_summary") or "").strip()
    actions = audit_json.get("action_items") or []
    actions_txt = "\n".join([f"- {a}" for a in actions if str(a).strip()])
    cc_comments = overall + (("\n\nAction items:\n" + actions_txt) if actions_txt else "")

    ctx = {
        # UNIT ticks
        "unit_cae": tick(unit.upper() == "CAE"),
        "unit_eece": tick(unit.upper() == "EECE"),
        "unit_mce": tick(unit.upper() == "MCE"),

        # Section 1 fields
        "specialization": specialization,
        "course_name": course_title,
        "academic_year": academic_year,
        "course_type_theory": "",       # set if you track this
        "course_type_practical": "",    # set if you track this
        "course_code": course_code,
        "semester": semester,
        "material_ppt": material_ppt,
        "material_reference": material_reference,
        "material_handout": material_handout,
        "lecturer_name": str((cdp_snapshot or {}).get("lecturer_name", "") or "").strip(),
        "date": time.strftime("%Y-%m-%d"),
        "material_lab_manual": material_lab_manual,

        # Section 2 Program Coordinator Review
        "pc_coverage_below": pc_cov_b, "pc_coverage_meet": pc_cov_m, "pc_coverage_above": pc_cov_a,
        "pc_coverage_remarks": remarks_blob(pc_cov),

        "pc_updated_below": pc_upd_b, "pc_updated_meet": pc_upd_m, "pc_updated_above": pc_upd_a,
        "pc_updated_remarks": remarks_blob(pc_upd),

        "pc_innov_below": pc_inn_b, "pc_innov_meet": pc_inn_m, "pc_innov_above": pc_inn_a,
        "pc_innov_remarks": remarks_blob(pc_inn),

        "pc_sequence_below": pc_seq_b, "pc_sequence_meet": pc_seq_m, "pc_sequence_above": pc_seq_a,
        "pc_sequence_remarks": remarks_blob(pc_seq),

        # Section 3 CC Review
        "cc_template_yes": cc_tpl_y, "cc_template_no": cc_tpl_n, "cc_template_remarks": remarks_blob(cc_tpl),
        "cc_outcomes_yes": cc_out_y, "cc_outcomes_no": cc_out_n, "cc_outcomes_remarks": remarks_blob(cc_out),
        "cc_numbering_yes": cc_num_y, "cc_numbering_no": cc_num_n, "cc_numbering_remarks": remarks_blob(cc_num),
        "cc_language_yes": cc_lng_y, "cc_language_no": cc_lng_n, "cc_language_remarks": remarks_blob(cc_lng),
        "cc_organization_yes": cc_org_y, "cc_organization_no": cc_org_n, "cc_organization_remarks": remarks_blob(cc_org),
        "cc_captions_yes": cc_cap_y, "cc_captions_no": cc_cap_n, "cc_captions_remarks": remarks_blob(cc_cap),
        "cc_citation_yes": cc_cit_y, "cc_citation_no": cc_cit_n, "cc_citation_remarks": remarks_blob(cc_cit),
        "cc_references_yes": cc_ref_y, "cc_references_no": cc_ref_n, "cc_references_remarks": remarks_blob(cc_ref),

        # Section 4 signatures/comments
        "pc_name_sign": pc_name_sign,
        "cc_comments": cc_comments,
        "cc_member_name_sign": cc_member_name_sign,
        "staff_ack_name_sign": staff_ack_name_sign,
    }

    tpl = DocxTemplate(str(template_path))
    tpl.render(ctx)
    buf = io.BytesIO()
    tpl.save(buf)
    return buf.getvalue()
